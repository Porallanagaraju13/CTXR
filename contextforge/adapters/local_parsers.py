import re
import fitz  # PyMuPDF
import docx
from pptx import Presentation
from pathlib import Path
from typing import List
from contextforge.core.models import DocumentPage, ExtractedTable

class LocalDocumentParser:
    """
    High-performance local extraction engines for PDF, Word (DOCX), and slides (PPTX).
    """

    @staticmethod
    def parse_pdf(filepath: Path) -> List[DocumentPage]:
        """
        Parses a PDF page-by-page using PyMuPDF.
        Includes high-performance layout-aware column-sorted extraction and table discovery.
        Reconstructs external links as markdown hyperlinks.
        """
        if not filepath.exists():
            raise FileNotFoundError(f"PDF file does not exist: {filepath}")
            
        doc = fitz.open(str(filepath))
        pages = []
        try:
            for page_num in range(len(doc)):
                page = doc[page_num]
                page_width = page.rect.width
                
                # Extract text blocks and sort them layout-aware (multi-column support)
                blocks = page.get_text("blocks")
                # Filter out image blocks (type 1 is image, type 0 is text in PyMuPDF)
                text_blocks = [b for b in blocks if b[6] == 0]
                
                # Column determination heuristic
                def get_block_column(b):
                    x0, y0, x1, y1 = b[0], b[1], b[2], b[3]
                    # Full-width titles/headers spanning > 70% of page width are column 0
                    if (x1 - x0) > page_width * 0.7:
                        return 0
                    # Right-aligned column
                    if x0 > page_width * 0.45:
                        return 1
                    # Left column
                    return 0
                
                # Sort blocks by column first, then by top-to-bottom y-coordinate
                sorted_blocks = sorted(text_blocks, key=lambda b: (get_block_column(b), b[1]))
                
                # Retrieve links on the current page
                links = page.get_links()
                
                processed_blocks = []
                for b in sorted_blocks:
                    block_rect = fitz.Rect(b[0], b[1], b[2], b[3])
                    block_text_clean = b[4].strip()
                    
                    # Intercept and insert hyperlinks in this block
                    for link in links:
                        if link.get("kind") == fitz.LINK_URI and "uri" in link:
                            link_rect = fitz.Rect(link["from"])
                            # Check if the link intersects with the block's rectangle
                            if (block_rect & link_rect).get_area() > 0:
                                link_text = page.get_text("text", clip=link_rect).strip()
                                if link_text:
                                    uri = link["uri"]
                                    markdown_link = f"[{link_text}]({uri})"
                                    if markdown_link not in block_text_clean:
                                        # Match exact or with whitespace variations
                                        if link_text in block_text_clean:
                                            block_text_clean = block_text_clean.replace(link_text, markdown_link, 1)
                                        else:
                                            escaped = re.escape(link_text)
                                            pattern = re.sub(r'(\\\s|\\n|\n)+', r'\\s+', escaped)
                                            try:
                                                block_text_clean = re.sub(pattern, markdown_link, block_text_clean, count=1)
                                            except Exception:
                                                pass
                                                
                    if block_text_clean:
                        processed_blocks.append(block_text_clean)
                        
                text = "\n\n".join(processed_blocks).strip()
                
                # Table discovery using native PyMuPDF structure finders
                extracted_tables = []
                try:
                    tables_found = page.find_tables()
                    for tab in tables_found:
                        data = tab.extract()
                        if data and len(data) > 0:
                            # Sanitize cell contents to avoid markdown layout breakages
                            headers = [str(h or "").strip().replace("\n", " ") for h in data[0]]
                            rows = []
                            for row in data[1:]:
                                rows.append([str(c or "").strip().replace("\n", " ") for c in row])
                            
                            # Construct GFM Markdown Table representation
                            md_repr = "| " + " | ".join(headers) + " |\n"
                            md_repr += "| " + " | ".join(["---"] * len(headers)) + " |\n"
                            for r in rows:
                                md_repr += "| " + " | ".join(r) + " |\n"
                                
                            extracted_tables.append(ExtractedTable(
                                headers=headers,
                                rows=rows,
                                markdown_representation=md_repr.strip()
                            ))
                except Exception as e:
                    # Table extraction errors should not crash the primary document text extraction
                    print(f"Non-fatal error parsing tables on page {page_num+1}: {e}")
                    
                pages.append(DocumentPage(
                     page_number=page_num + 1,
                     text_content=text,
                     tables=extracted_tables,
                     raw_size_bytes=len(text.encode("utf-8"))
                ))
        finally:
            doc.close()
            
        return pages

    @staticmethod
    def parse_docx(filepath: Path) -> List[DocumentPage]:
        """
        Parses a Word DOCX file, preserving structure, headers, lists, and tabular data.
        """
        if not filepath.exists():
            raise FileNotFoundError(f"DOCX file does not exist: {filepath}")
            
        doc = docx.Document(str(filepath))
        markdown_elements = []
        extracted_tables = []
        
        # Extract body text sequentially
        for p in doc.paragraphs:
            text = p.text.strip()
            if not text:
                continue
                
            style_name = p.style.name.lower()
            if "heading 1" in style_name:
                markdown_elements.append(f"# {text}")
            elif "heading 2" in style_name:
                markdown_elements.append(f"## {text}")
            elif "heading 3" in style_name:
                markdown_elements.append(f"### {text}")
            elif "list bullet" in style_name:
                markdown_elements.append(f"- {text}")
            elif "list number" in style_name:
                markdown_elements.append(f"1. {text}")
            else:
                markdown_elements.append(text)
                
        # Parse tables and build GFM blocks
        for table in doc.tables:
            if not table.rows:
                continue
            headers = [cell.text.strip().replace("\n", " ") for cell in table.rows[0].cells]
            rows = []
            for row in table.rows[1:]:
                rows.append([cell.text.strip().replace("\n", " ") for cell in row.cells])
                
            md_table = "| " + " | ".join(headers) + " |\n"
            md_table += "| " + " | ".join(["---"] * len(headers)) + " |\n"
            for r in rows:
                md_table += "| " + " | ".join(r) + " |\n"
                
            extracted_tables.append(ExtractedTable(
                headers=headers,
                rows=rows,
                markdown_representation=md_table.strip()
            ))
            
            # Embed table into main document stream
            markdown_elements.append(md_table.strip())
            
        full_text = "\n\n".join(markdown_elements)
        
        return [DocumentPage(
            page_number=1,
            text_content=full_text,
            tables=extracted_tables,
            raw_size_bytes=len(full_text.encode("utf-8"))
        )]

    @staticmethod
    def parse_pptx(filepath: Path) -> List[DocumentPage]:
        """
        Parses slides layout structures, drawing elements, text boxes, and slide notes.
        """
        if not filepath.exists():
            raise FileNotFoundError(f"PPTX file does not exist: {filepath}")
            
        prs = Presentation(str(filepath))
        pages = []
        
        for idx, slide in enumerate(prs.slides, start=1):
            slide_elements = [f"## Slide {idx}"]
            
            if slide.shapes.title:
                title_text = slide.shapes.title.text.strip()
                if title_text:
                    slide_elements.append(f"### {title_text}")
                    
            for shape in slide.shapes:
                if shape == slide.shapes.title:
                    continue
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if not text:
                            continue
                        level = paragraph.level
                        if level > 0:
                            slide_elements.append("  " * level + f"- {text}")
                        else:
                            slide_elements.append(text)
                            
            # Add speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_elements.append(f"\n> **Speaker Notes:** {notes}")
                    
            full_slide_text = "\n\n".join(slide_elements)
            pages.append(DocumentPage(
                page_number=idx,
                text_content=full_slide_text,
                raw_size_bytes=len(full_slide_text.encode("utf-8"))
            ))
            
        return pages
