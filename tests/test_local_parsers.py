import pytest
from pathlib import Path
import docx
from pptx import Presentation
from contextforge.adapters.local_parsers import LocalDocumentParser

def test_docx_parser_flow(tmp_path):
    doc_file = tmp_path / "document.docx"
    
    # Generate simple DOCX
    doc = docx.Document()
    doc.add_heading("ContextForge Architecture", level=1)
    p = doc.add_paragraph("High quality, low latency context pre-computation.")
    doc.save(str(doc_file))
    
    pages = LocalDocumentParser.parse_docx(doc_file)
    assert len(pages) == 1
    assert "ContextForge" in pages[0].text_content
    assert "# ContextForge Architecture" in pages[0].text_content

def test_pptx_parser_flow(tmp_path):
    prs_file = tmp_path / "presentation.pptx"
    
    # Generate simple PPTX
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "ContextForge Normalization"
    prs.save(str(prs_file))
    
    pages = LocalDocumentParser.parse_pptx(prs_file)
    assert len(pages) == 1
    assert "ContextForge Normalization" in pages[0].text_content

def test_pdf_parser_layout_and_hyperlinks(tmp_path):
    import fitz
    pdf_file = tmp_path / "test_doc.pdf"
    
    # Create a 2-column PDF document programmatically using fitz
    doc = fitz.open()
    page = doc.new_page(width=600, height=800)
    
    # Write column 2 text first, then column 1.
    # If the parser is layout-aware, it should sort column 1 before column 2.
    page.insert_textbox(fitz.Rect(400, 100, 550, 300), "This is column two text.\nAnd another line in column two.")
    page.insert_textbox(fitz.Rect(50, 100, 200, 300), "This is column one text.\nAnd another line in column one.")
    
    # Insert link in column 1
    rect = fitz.Rect(45, 85, 250, 150)
    page.insert_link({
        "kind": fitz.LINK_URI,
        "from": rect,
        "uri": "https://contextforge.ai"
    })
    
    doc.save(str(pdf_file))
    doc.close()
    
    # Parse the PDF
    pages = LocalDocumentParser.parse_pdf(pdf_file)
    assert len(pages) == 1
    content = pages[0].text_content
    
    # Verify that Column 1 text is sorted before Column 2 text
    # Column 1 has the link and should be formatted as a markdown link
    assert "https://contextforge.ai" in content
    assert "This is column two text." in content
    
    # Verify reading order: Column 1 should appear before Column 2
    idx1 = content.find("column one")
    idx2 = content.find("column two")
    assert idx1 < idx2
